"""PPTX connector for the OmniFetcher v1 canonical contract.

A PowerPoint deck is a natural composition tree, so this connector maps it
onto the canonical shape directly:

- a **deck** node (advisory ``kind`` ``"presentation"``) carrying the
  presentation-level descriptive fields (slide count, title, author) in its
  metadata core and namespaced ``source_extra["pptx"]``;
- one **slide** node child per slide (advisory ``kind`` ``"slide"``), each
  carrying its slide number / title in ``source_extra["pptx"]``;
- per-slide content as canonical atoms: a ``Text`` atom for the slide's text
  (title, body shapes, speaker notes), an ``Image`` atom per embedded picture,
  and a ``Table`` atom per embedded table.

Nodes are built only through :func:`omni_fetcher.v1.mapping.build_node`, so
the output contract stays uniform with every other source. The parse is
deterministic and read-only: bytes in, canonical tree out, no OCR and no
network beyond fetching the source bytes.

Failure handling follows the contract: a deck that cannot be opened at all is
a typed ``error(PARSE_ERROR)``; a deck that opens but has individual slides
that fail to extract yields a ``partial`` tree whose ``gaps`` name each
skipped slide, never a silent success with missing content.
"""

from __future__ import annotations

import asyncio
import io
from pathlib import Path
from typing import Any, AsyncIterator, Optional

import httpx

from omni_fetcher.v1.atoms import Image, Table, Text, TextFormat
from omni_fetcher.v1.auth import AuthCredential
from omni_fetcher.v1.errors import ErrorKind
from omni_fetcher.v1.fetcher import BaseFetcher
from omni_fetcher.v1.mapping import build_node
from omni_fetcher.v1.node import CompositionNode
from omni_fetcher.v1.result import (
    Gap,
    Result,
    from_exception,
    gap_from_exception,
    partial,
    success,
)
from omni_fetcher.v1.zoom import ZoomSpec

# Source namespace for descriptive PPTX fields placed in ``source_extra``.
_PPTX_NAMESPACE = "pptx"

# Advisory semantic kinds for the recursive deck -> slide composition tree.
_DECK_KIND = "presentation"
_SLIDE_KIND = "slide"

# MIME type of an Open-XML PowerPoint presentation.
_PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"


class PptxConnector(BaseFetcher):
    """
    PPTX connector for the v1 canonical contract
    ===============================================
    Maps a Microsoft PowerPoint presentation onto the canonical composition
    tree: a deck node (``kind="presentation"``) whose children are slide
    nodes (``kind="slide"``), each carrying ``Text`` / ``Image`` / ``Table``
    content atoms. It implements only the ``stream()`` primitive (``fetch()``
    is inherited from :class:`BaseFetcher`), reuses the v0.11 python-pptx
    extraction logic, and is deterministic and read-only (no OCR).
    ===============================================
    NOTE:
        1. Expected failures are returned as typed ``Result`` values, never
           raised: an unopenable deck is an ``error(PARSE_ERROR)`` and
           per-slide extraction failures are reported as ``partial`` gaps.
        2. Descriptive presentation/slide fields live in the metadata core
           and the namespaced ``source_extra["pptx"]`` mapping; atoms carry
           content only.
        3. Credentials are passed per call via ``auth`` and used transiently
           (only to fetch remote bytes); nothing is stored on the instance.

    Attributes
    ----------
        timeout:
            Per-call HTTP timeout, in seconds, for fetching remote decks.

    Methods
    -------
        can_handle:
        stream:
    """

    def __init__(self, timeout: float = 60.0) -> None:
        """
        Create a PPTX connector

        Parameters
        ----------
            timeout:
                Per-call HTTP timeout, in seconds, used when fetching a
                remote ``.pptx`` over HTTP(S).
        """
        self.timeout = timeout

    @classmethod
    def can_handle(cls, uri: str) -> bool:
        """
        Report whether this connector handles a URI

        Parameters
        ----------
            uri:
                The source URI to test.

        Return
        ------
            handled:
                ``True`` if ``uri`` names a ``.pptx`` resource.
        """
        return uri.lower().endswith(".pptx")

    async def stream(
        self,
        uri: str,
        *,
        auth: Optional[AuthCredential] = None,
        zoom: Optional[ZoomSpec] = None,
    ) -> AsyncIterator[Result]:
        """
        Stream the canonical result for a PPTX deck (the primitive)

        Fetches the deck bytes (local ``file://`` path or remote HTTP(S)
        URL), parses them with python-pptx off the event loop, and yields a
        single ``Result``: a ``success`` wrapping the canonical deck tree, a
        ``partial`` when one or more slides could not be fully extracted, or a
        typed ``error`` when the bytes cannot be fetched or opened at all.

        NOTE:
            1. A PPTX deck is a bounded source, so this stream yields exactly
               one terminal ``Result`` and then terminates; ``fetch()``
               collects it unchanged.
            2. ``zoom`` is accepted for contract conformance; the deck is
               emitted at its natural deck -> slide -> atom granularity.

        Parameters
        ----------
            uri:
                The ``.pptx`` source URI (``file://`` path or HTTP(S) URL).
            auth:
                The per-call credential, or ``None`` for unauthenticated
                sources. Used only to fetch remote bytes.
            zoom:
                Optional per-atom-type zoom spec; ``None`` means natural
                granularity.

        Return
        ------
            results:
                An async iterator yielding one terminal ``Result`` for the
                deck.
        """

        try:
            data = await self._load_bytes(uri, auth)
        except Exception as exc:  # noqa: BLE001 -- returned, never raised.
            yield from_exception(exc, locator=uri)
            return

        try:
            tree, gaps = await asyncio.to_thread(self._build_tree, data, uri)
        except Exception as exc:  # noqa: BLE001 -- returned, never raised.
            yield from_exception(
                exc,
                kind=ErrorKind.PARSE_ERROR,
                message="failed to parse PPTX deck",
                locator=uri,
            )
            return

        # Zoom is applied centrally (BaseFetcher.fetch / orchestrator.stream).
        # Slides already ARE the deck's sections, and central decomposition
        # preserves that: slide text is PLAIN, which has no section markers of
        # its own, so a SECTION request leaves the natural slide layer intact
        # (and does not gap). Finer levels decompose the slide text.
        result: Result = partial(tree, gaps) if gaps else success(tree)
        yield result

    async def _load_bytes(
        self,
        uri: str,
        auth: Optional[AuthCredential],
    ) -> bytes:
        """Fetch the raw deck bytes from a local path or remote URL."""
        if uri.startswith("file://"):
            return await asyncio.to_thread(self._read_local, uri)
        return await self._read_remote(uri, auth)

    @staticmethod
    def _read_local(uri: str) -> bytes:
        """Read a ``file://`` deck off the local filesystem."""
        path = uri[len("file://") :]
        if path.startswith("/") and len(path) > 2 and path[2] == ":":
            # Strip the leading slash from a Windows ``file:///C:/...`` path.
            path = path[1:]
        return Path(path).read_bytes()

    async def _read_remote(
        self,
        uri: str,
        auth: Optional[AuthCredential],
    ) -> bytes:
        """Fetch a remote deck over HTTP(S) using per-call credentials."""
        headers = _auth_headers(auth)
        async with httpx.AsyncClient(timeout=self.timeout) as client:
            response = await client.get(uri, headers=headers)
            response.raise_for_status()
            return response.content

    def _build_tree(
        self,
        data: bytes,
        uri: str,
    ) -> tuple[CompositionNode, list[Gap]]:
        """Parse deck bytes into the canonical deck -> slide -> atom tree."""
        from pptx import Presentation

        prs = Presentation(io.BytesIO(data))
        core = prs.core_properties
        title = core.title or None
        author = core.author or None
        slide_count = len(prs.slides)

        slide_nodes: list[CompositionNode] = []
        gaps: list[Gap] = []
        for number, slide in enumerate(prs.slides, start=1):
            locator = f"{uri}#slide={number}"
            try:
                node = self._build_slide_node(slide, number, uri)
            except Exception as exc:  # noqa: BLE001 -- captured as a gap.
                gaps.append(
                    gap_from_exception(
                        exc,
                        kind=ErrorKind.PARSE_ERROR,
                        locator=locator,
                    )
                )
                continue
            slide_nodes.append(node)

        deck = build_node(
            kind=_DECK_KIND,
            children=slide_nodes,
            author=author,
            source_url=uri,
            source_namespace=_PPTX_NAMESPACE,
            source_fields={
                "slide_count": slide_count,
                "title": title,
                "author": author,
                "mime_type": _PPTX_MIME,
            },
        )
        return deck, gaps

    def _build_slide_node(
        self,
        slide: Any,
        number: int,
        uri: str,
    ) -> CompositionNode:
        """Build one slide node with its Text / Image / Table atoms."""
        title = _slide_title(slide)
        text = _slide_text(slide, title)
        atoms: list[Any] = []
        if text:
            atoms.append(Text(content=text, format=TextFormat.PLAIN))
        atoms.extend(_slide_images(slide))
        atoms.extend(_slide_tables(slide))

        return build_node(
            kind=_SLIDE_KIND,
            atoms=atoms,
            source_url=f"{uri}#slide={number}",
            source_namespace=_PPTX_NAMESPACE,
            source_fields={
                "slide_number": number,
                "title": title,
            },
        )


def _auth_headers(auth: Optional[AuthCredential]) -> dict[str, str]:
    """Resolve a per-call credential into request headers (empty if none)."""
    if auth is None:
        return {}
    from omni_fetcher.v1.auth import NormalizedAuthResolver

    return NormalizedAuthResolver().resolve_headers(auth)


def _slide_title(slide: Any) -> Optional[str]:
    """Return the slide's title text, or ``None`` when it has none."""
    placeholder = slide.shapes.title
    if placeholder is None:
        return None
    text = placeholder.text
    return text.strip() if text else None


def _slide_text(slide: Any, title: Optional[str]) -> str:
    """Collect a slide's textual content (shapes + speaker notes)."""
    title_shape = slide.shapes.title
    parts: list[str] = []
    if title:
        parts.append(title)

    for shape in slide.shapes:
        if shape is title_shape:
            continue
        if getattr(shape, "has_text_frame", False) is False and not hasattr(shape, "text"):
            continue
        text = getattr(shape, "text", "")
        if text and text.strip():
            parts.append(text)

    notes = _speaker_notes(slide)
    if notes:
        parts.append(f"[Speaker Notes]: {notes}")

    return "\n".join(parts)


def _speaker_notes(slide: Any) -> Optional[str]:
    """Return the slide's speaker-notes text, or ``None`` when absent."""
    if not slide.has_notes_slide:
        return None
    notes_slide = slide.notes_slide
    frame = notes_slide.notes_text_frame
    if frame is None:
        return None
    text = frame.text
    return text.strip() if text and text.strip() else None


def _slide_images(slide: Any) -> list[Image]:
    """Extract embedded picture shapes as canonical ``Image`` atoms."""
    images: list[Image] = []
    for shape in slide.shapes:
        image = getattr(shape, "image", None)
        if image is None:
            continue
        try:
            blob = image.blob
            content_type = image.content_type
        except Exception:  # noqa: BLE001 -- skip a malformed picture.
            continue
        images.append(Image(format=content_type, data=blob))
    return images


def _slide_tables(slide: Any) -> list[Table]:
    """Extract embedded table shapes as canonical ``Table`` atoms."""
    tables: list[Table] = []
    for shape in slide.shapes:
        if getattr(shape, "has_table", False) is not True:
            continue
        grid = shape.table
        rows = [[cell.text for cell in row.cells] for row in grid.rows]
        tables.append(Table(rows=rows))
    return tables
