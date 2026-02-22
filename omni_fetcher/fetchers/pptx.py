"""PPTX fetcher for OmniFetcher."""

from __future__ import annotations

import io
from datetime import datetime
from pathlib import Path
from typing import Any, List, Optional

import httpx

try:
    import pptx
except ImportError:
    raise ImportError(
        "python-pptx is required for PPTX fetching. Install with: pip install omni_fetcher[office]"
    )

from omni_fetcher.core.registry import source
from omni_fetcher.fetchers.base import BaseFetcher
from omni_fetcher.schemas.atomics import (
    ImageDocument,
    TextDocument,
    TextFormat,
)
from omni_fetcher.schemas.base import DataCategory, FetchMetadata, MediaType
from omni_fetcher.schemas.documents import (
    PPTXDocument,
    SlideDocument,
)


@source(
    name="pptx",
    uri_patterns=["*.pptx"],
    mime_types=["application/vnd.openxmlformats-officedocument.presentationml.presentation"],
    priority=20,
    description="Fetch and parse Microsoft PowerPoint presentations",
)
class PPTXFetcher(BaseFetcher):
    """Fetcher for Microsoft PowerPoint presentations."""

    name = "pptx"
    priority = 20

    def __init__(self, timeout: float = 60.0):
        super().__init__()
        self.timeout = timeout

    @classmethod
    def can_handle(cls, uri: str) -> bool:
        """Check if this is a PPTX URL."""
        return uri.lower().endswith(".pptx")

    async def fetch(self, uri: str, **kwargs: Any) -> PPTXDocument:
        """Fetch and parse a PPTX document."""
        if uri.startswith("file://"):
            pptx_data = await self._fetch_local(uri)
        else:
            pptx_data = await self._fetch_remote(uri)

        metadata = FetchMetadata(
            source_uri=uri,
            fetched_at=datetime.now(),
            source_name=self.name,
            mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            file_size=len(pptx_data),
            fetch_duration_ms=None,
            cache_hit=False,
            status_code=None,
            encoding=None,
            last_modified=None,
            etag=None,
        )

        pptx_info = await self._parse_pptx(pptx_data, uri)

        return PPTXDocument(
            metadata=metadata,
            slides=pptx_info.get("slides", []),
            slide_count=pptx_info.get("slide_count", 0),
            author=pptx_info.get("author"),
            title=pptx_info.get("title"),
            category=DataCategory.DOCUMENT,
            media_type=MediaType.APPLICATION_PPTX,
        )

    async def _fetch_local(self, uri: str) -> bytes:
        """Fetch local PPTX file."""
        path = uri.replace("file://", "")
        if path.startswith("/"):
            path = path[1:]
        return Path(path).read_bytes()

    async def _fetch_remote(self, uri: str) -> bytes:
        """Fetch remote PPTX file."""
        headers = self.get_auth_headers()
        async with httpx.AsyncClient(timeout=self.timeout) as client:
            response = await client.get(uri, headers=headers)
            response.raise_for_status()
            return response.content

    async def _parse_pptx(self, pptx_data: bytes, uri: str) -> dict[str, Any]:
        """Parse PPTX and extract information."""

        def _parse():
            from pptx import Presentation

            prs = Presentation(io.BytesIO(pptx_data))

            info: dict[str, Any] = {
                "slides": [],
                "slide_count": len(prs.slides),
                "author": None,
                "title": None,
            }

            core_props = prs.core_properties
            info["title"] = core_props.title or None
            info["author"] = core_props.author or None

            for slide_num, slide in enumerate(prs.slides, start=1):
                slide_doc = self._extract_slide(slide, slide_num, uri)
                info["slides"].append(slide_doc)

            return info

        import asyncio

        loop = asyncio.get_event_loop()
        return await loop.run_in_executor(None, _parse)

    def _extract_slide(self, slide, slide_number: int, uri: str) -> SlideDocument:
        """Extract content from a single slide."""
        title: Optional[str] = None
        text_parts: List[str] = []
        images: List[ImageDocument] = []

        if slide.shapes.title:
            title = slide.shapes.title.text.strip() if slide.shapes.title.text else None

        for shape in slide.shapes:
            if shape == slide.shapes.title:
                continue

            if hasattr(shape, "text") and shape.text:
                text_parts.append(shape.text)

            if hasattr(shape, "image"):
                try:
                    image_blob = shape.image.blob
                    images.append(
                        ImageDocument(
                            source_uri=uri,
                            format=shape.image.content_type,
                            raw=image_blob,
                            file_name=f"slide{slide_number}_image{len(images)}",
                        )
                    )
                except Exception:
                    pass

        notes_slide = slide.notes_slide
        if notes_slide and notes_slide.notes_text_frame:
            notes_text = notes_slide.notes_text_frame.text.strip()
            if notes_text:
                text_parts.append(f"\n[Speaker Notes]: {notes_text}")

        full_text = "\n".join(text_parts)

        text_doc = TextDocument(
            source_uri=uri,
            content=full_text,
            format=TextFormat.PLAIN,
        )

        return SlideDocument(
            slide_number=slide_number,
            title=title,
            text=text_doc,
            images=images,
        )
